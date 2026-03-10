#!/usr/bin/env python3
"""
create-slides.py
Creates the "JDE Mission Control — Dashboard Guide" Google Slides presentation
as a .pptx file that can be uploaded to Google Slides.

Usage: python3 scripts/create-slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Output path ─────────────────────────────────────────────────────────────
OUTPUT = os.path.join(os.path.dirname(__file__), "..", "JDE_Mission_Control_Dashboard_Guide.pptx")

# ─── Colors ──────────────────────────────────────────────────────────────────
DARK_BG = RGBColor(17, 24, 39)
DARK_CARD = RGBColor(30, 41, 59)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(248, 250, 252)
BLUE = RGBColor(37, 99, 235)
LIGHT_BLUE = RGBColor(219, 234, 254)
DARK_BLUE = RGBColor(30, 58, 138)
GREEN = RGBColor(22, 163, 74)
LIGHT_GREEN = RGBColor(220, 252, 231)
DARK_GREEN = RGBColor(20, 83, 45)
ORANGE = RGBColor(245, 158, 11)
LIGHT_ORANGE = RGBColor(254, 243, 199)
DARK_ORANGE = RGBColor(120, 53, 15)
PURPLE = RGBColor(139, 92, 246)
LIGHT_PURPLE = RGBColor(237, 233, 254)
DARK_PURPLE = RGBColor(76, 29, 149)
RED = RGBColor(239, 68, 68)
LIGHT_RED = RGBColor(254, 226, 226)
DARK_RED = RGBColor(127, 29, 29)
GRAY = RGBColor(107, 114, 128)
LIGHT_GRAY = RGBColor(243, 244, 246)
DARK_TEXT = RGBColor(17, 24, 39)
MED_TEXT = RGBColor(75, 85, 99)


# ─── Helpers ─────────────────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=14, color=DARK_TEXT, bold=False, italic=False,
                alignment=PP_ALIGN.LEFT, line_spacing=None, font_name="Calibri",
                vertical_anchor=MSO_ANCHOR.TOP):
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Split by newlines and create paragraphs
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.font.name = font_name
        p.alignment = alignment
        if line_spacing:
            p.line_spacing = Pt(line_spacing)

    return txbox


def add_rounded_rect(slide, left, top, width, height, fill_color, corner_radius=0.15):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # No outline
    # Adjust corner radius
    try:
        shape.adjustments[0] = corner_radius
    except:
        pass
    return shape


# ─── Build presentation ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]  # Blank layout


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_textbox(slide, "🚀", 0, 0.8, 13.333, 1.2,
            font_size=72, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, "JDE Mission Control", 0, 2.0, 13.333, 1.0,
            font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, "Your Dashboard Guide", 0, 3.0, 13.333, 0.7,
            font_size=30, color=BLUE, bold=False, alignment=PP_ALIGN.CENTER)
add_textbox(slide, "Everything you need to know — explained simply.", 0, 4.2, 13.333, 0.5,
            font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2: What Is This?
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, OFF_WHITE)

add_textbox(slide, "🤔  What Is This Dashboard?", 0.8, 0.5, 12, 0.8,
            font_size=38, color=DARK_TEXT, bold=True)

add_rounded_rect(slide, 1.0, 1.7, 11.3, 4.5, WHITE)

add_textbox(slide, 'Think of it like a scoreboard at a sports game  🏟️', 1.5, 1.9, 10.3, 0.7,
            font_size=26, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide,
    "JDE Mission Control tracks everything that happens\n"
    "during a live car sales event at a dealership.\n"
    "\n"
    "📊  How many cars were sold\n"
    "💰  How much money was made\n"
    "👥  Who's selling the most\n"
    "🏆  Who's earning badges & awards",
    1.5, 2.8, 10.3, 2.8,
    font_size=20, color=MED_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=30)

add_textbox(slide, "Instead of whiteboards and spreadsheets → you have a live, real-time command center.",
            1.0, 6.3, 11.3, 0.5,
            font_size=16, color=GRAY, italic=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3: Quick Start
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_textbox(slide, "⚡  Quick Start — Do These 3 Things First", 0.8, 0.4, 12, 0.8,
            font_size=36, color=WHITE, bold=True)

# Box 1 - Blue
add_rounded_rect(slide, 0.6, 1.6, 3.7, 4.5, BLUE)
add_textbox(slide,
    "1️⃣\n\n"
    "Select Your\n"
    "Event\n\n"
    "Use the dropdown\n"
    "at the top of the\n"
    "sidebar to pick\n"
    "which event\n"
    "you're working on",
    0.6, 1.6, 3.7, 4.5,
    font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER, line_spacing=26)

# Box 2 - Green
add_rounded_rect(slide, 4.8, 1.6, 3.7, 4.5, GREEN)
add_textbox(slide,
    "2️⃣\n\n"
    "Log Your\n"
    "Deals\n\n"
    "Every time a car\n"
    "is sold, click\n"
    "\"New Deal\" and\n"
    "enter the info",
    4.8, 1.6, 3.7, 4.5,
    font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER, line_spacing=26)

# Box 3 - Orange
add_rounded_rect(slide, 9.0, 1.6, 3.7, 4.5, ORANGE)
add_textbox(slide,
    "3️⃣\n\n"
    "Check the\n"
    "Scoreboard\n\n"
    "Go to the\n"
    "Performance page\n"
    "to see charts,\n"
    "rankings & stats",
    9.0, 1.6, 3.7, 4.5,
    font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER, line_spacing=26)

add_textbox(slide, "That's it! Everything else updates automatically.  🎉", 0, 6.5, 13.333, 0.5,
            font_size=18, color=GRAY, italic=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4: Sidebar
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, OFF_WHITE)

add_textbox(slide, "🧭  The Sidebar — Your Map", 0.8, 0.5, 12, 0.8,
            font_size=38, color=DARK_TEXT, bold=True)

add_textbox(slide,
    "The sidebar is always on the left side of your screen.\n"
    "It's like the table of contents in a book — click any page to go there.",
    1.0, 1.3, 11.3, 0.8,
    font_size=18, color=MED_TEXT, line_spacing=26)

add_rounded_rect(slide, 1.0, 2.4, 11.3, 4.5, WHITE)

add_textbox(slide,
    "📊  Performance         →   Charts & leaderboard (the scoreboard)\n"
    "📝  Deals                      →   Every car sold (the deal log)\n"
    "🏆  Achievements        →   Badges, points & streaks\n"
    "📋  Daily Metrics          →   Enter daily numbers\n"
    "👥  Roster                      →   Your team of salespeople\n"
    "🚗  Inventory                →   Cars available at the dealership\n"
    "💵  Commissions          →   Who gets paid what\n"
    "⚙️  Settings                   →   App configuration",
    1.4, 2.6, 10.5, 4.0,
    font_size=19, color=MED_TEXT, line_spacing=32)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5: Event Switcher
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, OFF_WHITE)

add_textbox(slide, "🔄  The Event Switcher — Most Important Button!", 0.8, 0.5, 12, 0.8,
            font_size=34, color=DARK_TEXT, bold=True)

add_rounded_rect(slide, 1.0, 1.5, 11.3, 1.8, LIGHT_BLUE)
add_textbox(slide,
    "At the very top of the sidebar is a dropdown menu.\n"
    "Whichever event you pick here filters EVERYTHING across the entire app.",
    1.3, 1.7, 10.7, 1.4,
    font_size=20, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER, line_spacing=30)

add_textbox(slide, "💡  Think of it like switching TV channels", 1.0, 3.8, 11.3, 0.6,
            font_size=24, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide,
    "Each \"event\" = a different sales event at a dealership (usually 6 days).\n"
    "\n"
    "When you switch events:\n"
    "  •  All charts update to show that event's numbers\n"
    "  •  The leaderboard shows that event's team\n"
    "  •  Deals, roster, inventory — everything changes",
    1.3, 4.5, 10.7, 2.5,
    font_size=18, color=MED_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=28)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6: Performance Page Overview
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_textbox(slide, "📊  The Performance Page", 0.8, 0.4, 12, 0.8,
            font_size=38, color=WHITE, bold=True)
add_textbox(slide, "This is the heart of the dashboard — where you'll spend most of your time.",
            1.0, 1.2, 11.3, 0.5, font_size=18, color=GRAY, italic=True)

# KPI Cards box
add_rounded_rect(slide, 0.6, 2.0, 5.8, 2.0, DARK_CARD)
add_textbox(slide,
    "🔢  KPI Cards (Top)\n\n"
    "5 big number cards showing totals:\n"
    "Deals  •  Gross  •  Avg PVR  •  Close %  •  Ratios",
    0.9, 2.1, 5.2, 1.8,
    font_size=16, color=OFF_WHITE, line_spacing=24)

# Charts box
add_rounded_rect(slide, 6.9, 2.0, 5.8, 2.0, DARK_CARD)
add_textbox(slide,
    "📈  4 Charts (Middle)\n\n"
    "Gross per Day  •  Top Sellers\n"
    "Front vs Back  •  Daily Trend",
    7.2, 2.1, 5.2, 1.8,
    font_size=16, color=OFF_WHITE, line_spacing=24)

# Leaderboard box
add_rounded_rect(slide, 0.6, 4.3, 12.1, 2.0, DARK_CARD)
add_textbox(slide,
    "🏅  Leaderboard Table (Bottom)\n\n"
    "Every team member ranked by total gross — shows deals, ups, close %, gross breakdown, avg PVR, and earned badges",
    0.9, 4.4, 11.5, 1.8,
    font_size=16, color=OFF_WHITE, line_spacing=24)

add_textbox(slide, "⬆️  KPI Cards    →    📊  Charts    →    🏅  Leaderboard",
            0, 6.7, 13.333, 0.4,
            font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7: KPI Cards
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, OFF_WHITE)

add_textbox(slide, "🔢  The 5 Number Cards — Your Snapshot", 0.8, 0.4, 12, 0.8,
            font_size=36, color=DARK_TEXT, bold=True)
add_textbox(slide, "These sit at the very top of the Performance page. One glance = full picture.",
            1.0, 1.1, 11.3, 0.5, font_size=17, color=MED_TEXT)

# 5 mini cards
cards = [
    ("📦\nTotal\nDeals", "How many\ncars sold", LIGHT_BLUE, BLUE),
    ("💰\nTotal\nGross", "All profit\ncombined", LIGHT_GREEN, GREEN),
    ("📊\nAvg\nPVR", "Profit per\ncar sold", LIGHT_ORANGE, ORANGE),
    ("🎯\nClose\n%", "Deals ÷ Ups\n(walk-ins)", LIGHT_PURPLE, PURPLE),
    ("⚖️\nF:B\nRatio", "Front vs\nBack gross", LIGHT_RED, RED),
]

for i, (title, desc, bg_color, text_color) in enumerate(cards):
    x = 0.5 + i * 2.5
    add_rounded_rect(slide, x, 1.8, 2.2, 2.2, bg_color)
    add_textbox(slide, title, x, 1.9, 2.2, 1.2,
                font_size=16, color=text_color, bold=True, alignment=PP_ALIGN.CENTER, line_spacing=22)
    add_textbox(slide, desc, x, 3.1, 2.2, 0.8,
                font_size=13, color=text_color, alignment=PP_ALIGN.CENTER, line_spacing=20)

# Definitions box
add_rounded_rect(slide, 1.0, 4.4, 11.3, 2.6, WHITE)
add_textbox(slide,
    "🗣️  Quick Definitions:\n"
    "\n"
    "•  PVR = \"Per Vehicle Retailed\" — average profit per car\n"
    "•  Ups = customers who walked into the dealership\n"
    "•  Close % = what percentage of walk-ins actually bought a car\n"
    "•  Front Gross = profit from the car sale itself\n"
    "•  Back Gross = profit from financing, warranties, add-ons",
    1.3, 4.5, 10.7, 2.4,
    font_size=16, color=MED_TEXT, line_spacing=25)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8: Charts
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, OFF_WHITE)

add_textbox(slide, "📈  The 4 Charts — See the Trends", 0.8, 0.4, 12, 0.8,
            font_size=36, color=DARK_TEXT, bold=True)

charts = [
    (0.5, 1.5, "📊  Gross per Day", "Bar chart showing how much profit\nwas made each day of the event.\nLabels show \"5 sold • 184 ups\"", LIGHT_BLUE, BLUE),
    (6.9, 1.5, "🏆  Gross by Salesperson", "Horizontal bars ranking the top 10\nsellers by total profit.\nBiggest bar = top performer", LIGHT_GREEN, GREEN),
    (0.5, 4.2, "🍩  Front vs Back Breakdown", "Donut chart showing what % of\nprofit comes from the car sale\nvs. financing & add-ons", LIGHT_ORANGE, DARK_ORANGE),
    (6.9, 4.2, "📉  Daily PVR Trend", "Line chart tracking average profit\nper car over time. Going up = good!\nGoing down = adjust strategy", LIGHT_PURPLE, PURPLE),
]

for x, y, title, desc, bg_color, text_color in charts:
    add_rounded_rect(slide, x, y, 5.9, 2.4, bg_color)
    add_textbox(slide, title, x + 0.3, y + 0.2, 5.3, 0.5,
                font_size=20, color=text_color, bold=True)
    add_textbox(slide, desc, x + 0.3, y + 0.8, 5.3, 1.4,
                font_size=16, color=text_color, line_spacing=24)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9: Leaderboard
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_textbox(slide, "🏅  The Leaderboard — Who's Winning?", 0.8, 0.4, 12, 0.8,
            font_size=36, color=WHITE, bold=True)
add_textbox(slide, "A table ranking every team member from highest to lowest total gross profit.",
            1.0, 1.15, 11.3, 0.5, font_size=18, color=GRAY)

# Table header
add_rounded_rect(slide, 0.6, 1.9, 12.1, 0.6, DARK_CARD)
add_textbox(slide, "#      Name                    Role             Deals      Ups      Close%      Front        Back        Total         PVR        Badges",
            0.8, 1.95, 11.7, 0.5, font_size=14, color=GRAY, bold=True)

# Row 1
add_rounded_rect(slide, 0.6, 2.6, 12.1, 0.55, RGBColor(30, 50, 70))
add_textbox(slide, "1      John Smith           Sales              12          40        30%        $24K        $18K        $42K       $3.5K      🎯🔥",
            0.8, 2.65, 11.7, 0.45, font_size=14, color=OFF_WHITE)

# Row 2
add_rounded_rect(slide, 0.6, 3.2, 12.1, 0.55, RGBColor(25, 38, 55))
add_textbox(slide, "2      Jane Doe              Sales               8           35        23%        $16K        $12K        $28K       $3.5K      🏆",
            0.8, 3.25, 11.7, 0.45, font_size=14, color=OFF_WHITE)

# Explanation
add_rounded_rect(slide, 0.6, 4.2, 12.1, 2.6, DARK_CARD)
add_textbox(slide,
    "What each column means:\n"
    "\n"
    "#  =  Rank (1st = most profit)              Deals  =  Cars sold              Ups  =  Customers seen\n"
    "Close%  =  How many ups became sales      Front  =  Car sale profit      Back  =  F&I profit\n"
    "Total  =  Front + Back combined               PVR  =  Avg profit per car    Badges  =  Awards earned",
    0.9, 4.3, 11.5, 2.4,
    font_size=16, color=OFF_WHITE, line_spacing=28)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10: Deals Page
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, OFF_WHITE)

add_textbox(slide, "📝  The Deals Page — Every Car Sold", 0.8, 0.5, 12, 0.8,
            font_size=36, color=DARK_TEXT, bold=True)
add_textbox(slide, "This is your deal log — a list of every single car sale at the event.",
            1.0, 1.2, 11.3, 0.5, font_size=18, color=MED_TEXT)

add_rounded_rect(slide, 1.0, 2.0, 11.3, 1.6, LIGHT_GREEN)
add_textbox(slide,
    "Each row = one car sale, with info like:\n"
    "🚗 Vehicle (stock #, year, make, model)  •  👤 Customer  •  🧑‍💼 Salesperson\n"
    "💰 Front Gross  •  💵 Back Gross  •  📊 Total Gross  •  📅 Date",
    1.3, 2.1, 10.7, 1.4,
    font_size=18, color=GREEN, alignment=PP_ALIGN.CENTER, line_spacing=28)

add_textbox(slide, "How to log a new deal:", 1.0, 4.0, 11.3, 0.5,
            font_size=24, color=DARK_TEXT, bold=True)

add_textbox(slide,
    "1.  Click the \"New Deal\" button in the top right\n"
    "2.  Fill in the vehicle, customer, and gross numbers\n"
    "3.  Pick the salesperson from the roster dropdown\n"
    "4.  Hit Save — done! ✅",
    1.3, 4.6, 10.7, 2.2,
    font_size=20, color=MED_TEXT, line_spacing=32)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11: What Happens When You Log a Deal
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_textbox(slide, "⚡  What Happens When You Log a Deal?", 0.8, 0.3, 12, 0.8,
            font_size=36, color=WHITE, bold=True)
add_textbox(slide, "A lot of magic happens behind the scenes — automatically!",
            1.0, 1.1, 11.3, 0.5, font_size=17, color=GRAY, italic=True)

# 5 chain boxes - row 1 (3 boxes)
colors1 = [BLUE, GREEN, ORANGE]
texts1 = [
    "1️⃣\n\n💾  Deal Saved\n\nThe sale is recorded\nin the database",
    "2️⃣\n\n🏆  Badges Check\n\nDid this earn any\nnew achievements?",
    "3️⃣\n\n🔥  Streak Updated\n\nConsecutive days\nwith a sale tracked",
]
for i, (color, text) in enumerate(zip(colors1, texts1)):
    x = 0.5 + i * 4.25
    add_rounded_rect(slide, x, 1.9, 3.7, 2.3, color)
    add_textbox(slide, text, x, 2.0, 3.7, 2.2,
                font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER, line_spacing=24)

# Row 2 (2 boxes)
colors2 = [PURPLE, RED]
texts2 = [
    "4️⃣\n\n🚗  Inventory Synced\n\nVehicle marked as \"sold\"",
    "5️⃣\n\n🎉  Toast Pops Up!\n\n\"Badge Earned: First Blood!\"",
]
for i, (color, text) in enumerate(zip(colors2, texts2)):
    x = 2.6 + i * 4.4
    add_rounded_rect(slide, x, 4.5, 3.7, 2.2, color)
    add_textbox(slide, text, x, 4.6, 3.7, 2.1,
                font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER, line_spacing=24)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12: Achievements Overview
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, OFF_WHITE)

add_textbox(slide, "🏆  Achievements — Badges, Points & Streaks", 0.8, 0.4, 12, 0.8,
            font_size=34, color=DARK_TEXT, bold=True)

add_textbox(slide,
    "Think of it like earning trophies in a video game  🎮\n"
    "Sell more, sell better → earn badges → climb the points leaderboard",
    1.0, 1.2, 11.3, 0.8,
    font_size=19, color=MED_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=28)

# 3 tab cards
tabs = [
    ("Tab 1: Badges  🎖️",
     "All 18 badges shown\nin a grid.\n\n✅ Earned = full color\n🔒 Locked = grayed out\n\nEach badge has\na point value",
     LIGHT_BLUE, BLUE),
    ("Tab 2: Team  👥",
     "Table showing each\nteam member's:\n\n•  Badges earned\n•  Total points\n•  Current streak\n•  Recent badge",
     LIGHT_GREEN, GREEN),
    ("Tab 3: Points  🥇",
     "Leaderboard ranked\nby total badge points.\n\nSeparate from the\ngross leaderboard —\nthis is about\nachievements!",
     LIGHT_ORANGE, DARK_ORANGE),
]

for i, (title, desc, bg_color, text_color) in enumerate(tabs):
    x = 0.5 + i * 4.25
    add_rounded_rect(slide, x, 2.5, 3.8, 4.2, bg_color)
    add_textbox(slide, title, x, 2.7, 3.8, 0.5,
                font_size=20, color=text_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, desc, x, 3.3, 3.8, 3.2,
                font_size=16, color=text_color, alignment=PP_ALIGN.CENTER, line_spacing=24)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13: Badge Categories
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_textbox(slide, "🎖️  18 Badges Across 5 Categories", 0.8, 0.3, 12, 0.8,
            font_size=36, color=WHITE, bold=True)

# Row 1 - 3 categories
cats1 = [
    ("🔵  Sales (5 badges)", "First Blood  •  Hat Trick\n5-Pack  •  10-Unit Club\n15-Car Legend", DARK_BLUE, LIGHT_BLUE),
    ("🟢  Gross (4 badges)", "$10K Day  •  $25K Day\n$50K Total\n$100K Club", DARK_GREEN, LIGHT_GREEN),
    ("🟣  Closing (3 badges)", "Sharpshooter (20%+)\nSniper (30%+)\nCloser Supreme (40%+)", DARK_PURPLE, LIGHT_PURPLE),
]

for i, (title, desc, bg_color, text_color) in enumerate(cats1):
    x = 0.4 + i * 4.3
    add_rounded_rect(slide, x, 1.3, 3.8, 2.2, bg_color)
    add_textbox(slide, title, x, 1.4, 3.8, 0.5,
                font_size=17, color=text_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, desc, x, 2.1, 3.8, 1.2,
                font_size=15, color=text_color, alignment=PP_ALIGN.CENTER, line_spacing=22)

# Row 2 - 2 categories
cats2 = [
    ("🟠  Streak (3 badges)", "On a Roll (2 days)\nHot Streak (3 days)\nIron Man (5+ days)", DARK_ORANGE, LIGHT_ORANGE),
    ("🔴  Team (3 badges)", "Top Dog (#1 on board)\nComeback Kid\nClean Sheet (0 washouts)", DARK_RED, LIGHT_RED),
]

for i, (title, desc, bg_color, text_color) in enumerate(cats2):
    x = 2.5 + i * 4.5
    add_rounded_rect(slide, x, 3.8, 3.8, 2.2, bg_color)
    add_textbox(slide, title, x, 3.9, 3.8, 0.5,
                font_size=17, color=text_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, desc, x, 4.6, 3.8, 1.2,
                font_size=15, color=text_color, alignment=PP_ALIGN.CENTER, line_spacing=22)

add_textbox(slide, "Badges are earned automatically when you hit the target — no manual action needed!",
            0, 6.4, 13.333, 0.5,
            font_size=16, color=GRAY, italic=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14: Streaks
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, OFF_WHITE)

add_textbox(slide, "🔥  Streaks — Keep the Momentum Going", 0.8, 0.4, 12, 0.8,
            font_size=36, color=DARK_TEXT, bold=True)

add_rounded_rect(slide, 1.0, 1.5, 11.3, 2.0, LIGHT_ORANGE)
add_textbox(slide,
    "A streak counts how many days IN A ROW you've made at least one sale.\n"
    "Sell today + sell tomorrow = 2-day streak  🔥🔥\n"
    "Miss a day? Streak resets back to 1.",
    1.3, 1.7, 10.7, 1.6,
    font_size=20, color=DARK_ORANGE, alignment=PP_ALIGN.CENTER, line_spacing=30)

add_textbox(slide, "How it shows up:", 1.0, 3.8, 11.3, 0.5,
            font_size=24, color=DARK_TEXT, bold=True)

add_rounded_rect(slide, 1.0, 4.5, 11.3, 2.3, WHITE)
add_textbox(slide,
    "🔥 3   (best: 5)     ← This means:\n"
    "\n"
    "•  Current streak: 3 consecutive days with a sale\n"
    "•  Best ever streak: 5 days (their personal record)\n"
    "•  🟠 Orange flame = 3+ days   •   🟡 Yellow = 1-2 days   •   ⚪ Gray = no streak",
    1.3, 4.6, 10.7, 2.0,
    font_size=17, color=MED_TEXT, line_spacing=28)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15: Daily Metrics
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, OFF_WHITE)

add_textbox(slide, "📋  Daily Metrics — Enter Your Daily Numbers", 0.8, 0.4, 12, 0.8,
            font_size=34, color=DARK_TEXT, bold=True)
add_textbox(slide, "A simple spreadsheet where you type in the daily totals for the event.",
            1.0, 1.2, 11.3, 0.5, font_size=18, color=MED_TEXT)

add_rounded_rect(slide, 1.0, 1.9, 11.3, 1.8, WHITE)
add_textbox(slide,
    "Each row = one day of the event:\n"
    "\n"
    "📅 Date   •   👣 Ups (walk-ins)   •   🚗 Sold   •   💰 Total Gross\n"
    "💵 Front Gross   •   🏦 Back Gross   •   📝 Notes",
    1.3, 2.0, 10.7, 1.6,
    font_size=18, color=MED_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=28)

add_textbox(slide, "How to use it:", 1.0, 4.0, 11.3, 0.5,
            font_size=24, color=DARK_TEXT, bold=True)

add_textbox(slide,
    "1.  Click \"Add Day\" to add a new row (date auto-fills)\n"
    "2.  Type numbers directly into the cells\n"
    "3.  Changed rows turn yellow so you know what needs saving\n"
    "4.  Click \"Save Changes\" to save everything at once\n"
    "5.  Close % auto-calculates (you don't type it)",
    1.3, 4.6, 10.7, 2.5,
    font_size=19, color=MED_TEXT, line_spacing=32)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16: Roster & Inventory
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_textbox(slide, "👥  Roster  &  🚗  Inventory", 0.8, 0.4, 12, 0.8,
            font_size=38, color=WHITE, bold=True)

# Roster box
add_rounded_rect(slide, 0.6, 1.5, 5.7, 5.0, DARK_CARD)
add_textbox(slide,
    "👥  Roster\n"
    "\n"
    "Your team for this event.\n"
    "\n"
    "Each person has:\n"
    "  •  Name\n"
    "  •  Role (Sales, Closer,\n"
    "     Team Leader, F&I)\n"
    "\n"
    "When you log a deal,\n"
    "you pick the salesperson\n"
    "from this list.\n"
    "\n"
    "Everyone on the roster\n"
    "shows up on the leaderboard.",
    0.9, 1.6, 5.1, 4.8,
    font_size=17, color=OFF_WHITE, line_spacing=24)

# Inventory box
add_rounded_rect(slide, 7.0, 1.5, 5.7, 5.0, DARK_CARD)
add_textbox(slide,
    "🚗  Inventory\n"
    "\n"
    "Cars available at the\n"
    "dealership for this event.\n"
    "\n"
    "Each vehicle has:\n"
    "  •  Stock number\n"
    "  •  Year, Make, Model\n"
    "  •  Status (available/sold)\n"
    "\n"
    "When you log a deal\n"
    "with a stock number,\n"
    "that car automatically\n"
    "gets marked as \"sold\".",
    7.3, 1.6, 5.1, 4.8,
    font_size=17, color=OFF_WHITE, line_spacing=24)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17: Cheat Sheet
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, OFF_WHITE)

add_textbox(slide, "📌  Quick Reference Cheat Sheet", 0.8, 0.3, 12, 0.7,
            font_size=36, color=DARK_TEXT, bold=True)

add_rounded_rect(slide, 0.5, 1.2, 12.3, 5.8, WHITE)

add_textbox(slide,
    "\"I want to...\"                                                                              →  Go here\n"
    "\n"
    "🔄  Switch which event I'm looking at                   →  Event Switcher (top of sidebar)\n"
    "📊  See charts and rankings                                    →  Performance page\n"
    "📝  Log a car sale                                                     →  Deals  →  New Deal button\n"
    "🏆  See who earned badges                                     →  Achievements page\n"
    "📋  Enter daily ups / sold / gross numbers             →  Daily Metrics page\n"
    "👥  Add a new team member                                   →  Roster page\n"
    "🚗  Check what cars are left                                     →  Inventory page\n"
    "💵  See who gets paid what                                     →  Commissions page\n"
    "✏️  Edit a deal                                                            →  Deals  →  click the deal row\n"
    "🔄  Refresh the data                                                  →  Click \"Refresh\" on Performance",
    0.9, 1.4, 11.5, 5.4,
    font_size=17, color=MED_TEXT, line_spacing=30)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18: Closing
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_textbox(slide, "🎉", 0, 0.8, 13.333, 1.2,
            font_size=72, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, "You're Ready!", 0, 2.1, 13.333, 1.0,
            font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, "Remember the 3 steps:", 0, 3.3, 13.333, 0.6,
            font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide,
    "1.  Pick your event   🔄\n"
    "2.  Log your deals   📝\n"
    "3.  Watch the scoreboard   📊",
    0, 4.0, 13.333, 1.8,
    font_size=28, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER, line_spacing=38)

add_textbox(slide, "Everything else happens automatically. Go sell some cars!  🚗💨",
            0, 6.0, 13.333, 0.6,
            font_size=20, color=GRAY, italic=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"✅ Presentation saved to: {os.path.abspath(OUTPUT)}")
print(f"   {len(prs.slides)} slides created")
print(f"\n   To use in Google Slides:")
print(f"   1. Go to slides.google.com")
print(f"   2. Click the blank presentation (or open the one already created)")
print(f"   3. File → Import slides → Upload → select this .pptx file")
print(f"   4. Select All → Import")
